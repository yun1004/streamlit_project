import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import io
import tempfile
from PIL import Image

st.title("PDF to PPTX Converter 🔄")

uploaded_file = st.file_uploader("PDF 파일을 업로드하세요", type="pdf")

if uploaded_file is not None:
    st.write("파일 변환 중...")
    
    # PDF를 이미지로 변환
    try:
        images = convert_from_bytes(uploaded_file.getvalue())
    except Exception as e:
        st.error(f"PDF 변환 중 오류 발생: {str(e)}")
        st.stop()
    
    # 새 PowerPoint 프레젠테이션 생성
    prs = Presentation()
    
    # PDF의 첫 페이지 크기로 슬라이드 크기 설정
    prs.slide_width = Inches(images[0].width / 96)
    prs.slide_height = Inches(images[0].height / 96)
    
    for image in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드 레이아웃
        
        # 이미지를 슬라이드에 추가
        left = top = Inches(0)
        
        # 임시 파일을 사용하여 이미지 저장
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
            image.save(temp_file, format='PNG')
            temp_file_path = temp_file.name
        
        try:
            slide.shapes.add_picture(temp_file_path, left, top, 
                                     width=prs.slide_width, 
                                     height=prs.slide_height)
        except Exception as e:
            st.error(f"이미지 추가 중 오류 발생: {str(e)}")
            continue
    
    # PPTX 파일 저장
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    
    st.success("변환 완료!")
    
    # 다운로드 버튼 생성
    st.download_button(
        label="PPTX 파일 다운로드",
        data=pptx_file,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
